import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def replace_text_preserve_format(shape, new_text, width_inches=None, force_color=None, force_font_size=None):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    
    # Save original style properties
    orig_font_name = "Arial"
    orig_font_size = Pt(13) # Enforce 13pt as default for readability
    orig_font_color = RGBColor(51, 51, 51) # Default dark gray
    orig_bold = False
    
    if tf.paragraphs and tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        if run.font.name:
            orig_font_name = run.font.name
        if run.font.size:
            orig_font_size = run.font.size
        if run.font.color and run.font.color.type != 0:
            try:
                orig_font_color = run.font.color.rgb
            except Exception:
                pass
        orig_bold = run.font.bold

    if force_color:
        orig_font_color = force_color
    if force_font_size:
        orig_font_size = Pt(force_font_size)

    # Clear all text
    tf.text = ""
    
    lines = new_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.name = orig_font_name
        run.font.size = orig_font_size
        run.font.color.rgb = orig_font_color
        run.font.bold = orig_bold

    # Adjust width if requested to prevent overlapping
    if width_inches:
        shape.width = Inches(width_inches)

def add_styled_textbox(slide, left, top, width, height, text, font_size, is_bold=False, font_name="Arial", color=RGBColor(51, 51, 51)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = is_bold
    return txBox

def main():
    template_path = r"E:\Downloads\Telegram Desktop\Норникель 2026.pptx"
    output_path = r"E:\Downloads\Telegram Desktop\GraphAnalyzer_Презентация.pptx"
    
    prs = Presentation(template_path)
    
    # ── Слайд 1: Титульный ──
    slide1 = prs.slides[0]
    # Details text box below the capsule
    add_styled_textbox(slide1, Inches(2.0), Inches(4.7), Inches(9.0), Inches(2.0),
        "Решение трека «Научный клубок» — Система GraphAnalyzer\n"
        "Команда «Четыре туза»\n"
        "Состав: Минниханов Д., Кубинский А., Селезнев М., Корниенко К.\n"
        "Веб-стенд: danil.goida2.ru/graph  ·  GitHub: github.com/DanilMini1/graph-visualizer",
        13, is_bold=True, color=RGBColor(111, 44, 232)
    )
    
    # ── Слайд 2: Методология ──
    slide2 = prs.slides[1]
    slide2.shapes[0].width = Inches(9.0) # Extend category width to prevent cut-off
    slide2.shapes[1].width = Inches(9.0) # Extend title width
    replace_text_preserve_format(slide2.shapes[0], "01 / МЕТОДОЛОГИЯ")
    replace_text_preserve_format(slide2.shapes[1], "Графовый R&D анализ научного опыта материаловедения")
    replace_text_preserve_format(slide2.shapes[2], 
        "• ОБЪЕКТ: Массивы неструктурированных научных R&D публикаций и R&D патентов.\n"
        "• ПРЕДМЕТ: AI-алгоритмы извлечения знаний и инструменты анализа связей.\n"
        "• ЦЕЛЬ: Разработка GraphAnalyzer — системы поиска скрытых R&D связей.\n"
        "• ЗАДАЧИ:\n"
        "  1. Проектирование открытого формата обмена данными graph.json.\n"
        "  2. Разработка D3-Force layout для автоматической R&D кластеризации.\n"
        "  3. Реализация клиентского естественно-языкового NLP-поисковика.\n"
        "  4. Разработка алгоритмов поиска путей для связывания R&D отчетов.",
        width_inches=5.3, force_font_size=13
    )
    
    # ── Слайд 3: Проблема и Бизнес-ценность ──
    slide3 = prs.slides[2]
    slide3.shapes[0].width = Inches(9.0)
    slide3.shapes[1].width = Inches(9.0)
    replace_text_preserve_format(slide3.shapes[0], "02 / БИЗНЕС-КОНТЕКСТ")
    replace_text_preserve_format(slide3.shapes[1], "Тяжелые AI-решения дороги, а R&D данные нечитаемы")
    replace_text_preserve_format(slide3.shapes[2], 
        "• ВЗРЫВ ИНФОРМАЦИИ: Рост числа публикаций. Ученые тратят до 40% времени на поиск.\n"
        "• СКРЫТЫЕ СВЯЗИ: Зависимости R&D данных остаются незамеченными.\n"
        "• НАШЕ РЕШЕНИЕ: Разделение парсинга (внешняя/локальная LLM делает graph.json) и визуального анализа (GraphAnalyzer на клиенте).\n"
        "• R&D ЭФФЕКТ: Снижение времени поиска на 70%, 100% конфиденциальность данных.",
        width_inches=5.3, force_font_size=13
    )
    
    # ── Слайд 4: Архитектура решения и Стек ──
    slide4 = prs.slides[3]
    slide4.shapes[0].width = Inches(9.0)
    slide4.shapes[1].width = Inches(9.0)
    replace_text_preserve_format(slide4.shapes[0], "03 / ТЕХНОЛОГИИ")
    replace_text_preserve_format(slide4.shapes[1], "Next.js 16 + React Flow (xyflow) + D3.js для больших графов")
    replace_text_preserve_format(slide4.shapes[2], 
        "• ЯДРО: Next.js 16 (App Router) и React 19 для максимальной производительности.\n"
        "• ВИЗУАЛИЗАЦИЯ: React Flow (xyflow) с аппаратным ускорением для плавного зума.\n"
        "• РАСКЛАДКА: d3-force для физической симуляции сил с учетом центральности.\n"
        "• ДЕПЛОЙ: Standalone Next.js Docker-контейнер на ВМ под прокси Caddy.\n"
        "• ДОСТУПНОСТЬ: Стенд полностью доступен по адресу danil.goida2.ru/graph.",
        width_inches=5.3, force_font_size=13
    )
    
    # ── Слайд 5: Ключевая фича 1 — Локальный NLP Engine ──
    slide5 = prs.slides[4]
    slide5.shapes[1].width = Inches(9.0)
    slide5.shapes[2].width = Inches(9.0)
    replace_text_preserve_format(slide5.shapes[1], "04 / ФУНКЦИОНАЛ")
    replace_text_preserve_format(slide5.shapes[2], "Клиентский семантический поиск с синонимами металлургии")
    replace_text_preserve_format(slide5.shapes[3], 
        "• СЛОВАРЬ СИНОНИМОВ: База из 30+ R&D групп терминов (электроэкстракция ↔ electrowinning).\n"
        "• ПАРСЕР ЗАПРОСОВ: Выделение географии (Россия/зарубеж), дат («5 лет») и числовых диапазонов параметров («Ni > 5%»).\n"
        "• РАНЖИРОВАНИЕ: Сортировка по совпадениям в названиях с учетом весов важности.",
        width_inches=5.3, force_font_size=13
    )
    # Clear dummy shapes and capsule text buttons on slide 5
    replace_text_preserve_format(slide5.shapes[4], "")
    replace_text_preserve_format(slide5.shapes[5], "")
    
    # ── Слайд 6: Ключевая фича 2 — Трассировщик связей (Две карточки) ──
    slide6 = prs.slides[5]
    # We populate the text INSIDE the two cards of slide 6
    replace_text_preserve_format(slide6.shapes[0], 
        "05 / ТРАССИРОВКА\n\n"
        "Нахождение цепочек связей:\n\n"
        "• Выявление скрытых зависимостей между научными R&D исследованиями.\n"
        "• Поиск кратчайшего R&D пути (алгоритмы Дейкстры и BFS) по ребрам графа.\n"
        "• Анимированное отображение связей и переходов.",
        force_color=RGBColor(51, 51, 51), force_font_size=12
    )
    replace_text_preserve_format(slide6.shapes[1], 
        "05 / БИЗНЕС-ПРИМЕНЕНИЕ\n\n"
        "Интерактивная аналитика:\n\n"
        "• Ученый-металлург видит пошаговую логическую цепочку переходов между процессами, материалами и свойствами.\n"
        "• Отображение метаданных R&D отчетов для каждого узла на пути.\n"
        "• Экономия R&D времени на R&D литературу.",
        force_color=RGBColor(51, 51, 51), force_font_size=12
    )
    
    # ── Слайд 7: Демонстрация работы (Интерфейс) ──
    slide7 = prs.slides[6]
    add_styled_textbox(slide7, Inches(0.6), Inches(1.19), Inches(5.5), Inches(0.8), "06 / ДЕМОНСТРАЦИЯ", 12, is_bold=True, color=RGBColor(111, 44, 232))
    add_styled_textbox(slide7, Inches(0.6), Inches(1.83), Inches(12.0), Inches(0.8), "Интерактивный дашборд GraphAnalyzer в действии", 20, is_bold=True, color=RGBColor(30, 20, 80))
    add_styled_textbox(slide7, Inches(0.6), Inches(2.90), Inches(6.0), Inches(4.0),
        "• ИНТЕРФЕЙС: Левая панель содержит вкладки «Навигация» и «Маршруты».\n"
        "• ХОЛСТ: Правая панель отображает граф. Узлы окрашены по сообществам и масштабируются по важности.\n"
        "• ИНФО-ПАНЕЛИ: Подробная информация по выбранному узлу и результаты трассировки путей выходят в R&D карточки.\n"
        "• ЖИВОЙ СТЕНД: Приложение полностью доступно для R&D тестирования по адресу danil.goida2.ru/graph.",
        13, color=RGBColor(51, 51, 51)
    )
    
    # ── Слайд 8: Масштабируемость ──
    slide8 = prs.slides[7]
    add_styled_textbox(slide8, Inches(0.6), Inches(1.19), Inches(5.5), Inches(0.8), "07 / ЭФФЕКТИВНОСТЬ", 12, is_bold=True, color=RGBColor(111, 44, 232))
    add_styled_textbox(slide8, Inches(0.6), Inches(1.83), Inches(12.0), Inches(0.8), "Масштабируемость и бизнес-эффект для Норникеля", 20, is_bold=True, color=RGBColor(30, 20, 80))
    add_styled_textbox(slide8, Inches(0.6), Inches(2.90), Inches(7.5), Inches(4.0),
        "• УНИВЕРСАЛЬНОСТЬ: Интеграция с любыми LLM через простой API-контракт graph.json.\n"
        "• БЕЗОПАСНОСТЬ: Полный R&D контроль над данными. Разбор запросов и поиск происходят локально, без внешних серверов.\n"
        "• МАСШТАБИРУЕМОСТЬ: Легкое расширение словаря синонимов под новые цеха, экологические или геологические задачи.\n"
        "• ЭФФЕКТ: Снижение трудоемкости анализа литературы на 70%, исключение дублирования исследований и ошибочных экспериментов.",
        13, color=RGBColor(51, 51, 51)
    )
    
    # ── Слайд 9: Команда и итоги ──
    slide9 = prs.slides[8]
    add_styled_textbox(slide9, Inches(0.6), Inches(1.19), Inches(5.5), Inches(0.8), "08 / КОМАНДА", 12, is_bold=True, color=RGBColor(111, 44, 232))
    add_styled_textbox(slide9, Inches(0.6), Inches(1.83), Inches(12.0), Inches(0.8), "Команда «Четыре туза»: Синергия системного анализа, ML и R&D", 20, is_bold=True, color=RGBColor(30, 20, 80))
    add_styled_textbox(slide9, Inches(0.6), Inches(2.90), Inches(8.0), Inches(4.0),
        "• Минниханов Данил — Системный анализ, серверный деплой, Caddy-проксирование, UI-дизайн и верстка презентации.\n"
        "• Кубинский Арсений — R&D-аналитика металлургических R&D процессов, сбор и подготовка R&D наборов данных.\n"
        "• Михаил Селезнев — ML-разработка, d3-layout, клиентский NLP Engine, реализация графовых путей Дейкстры.\n"
        "• Корниенко Константин — Frontend-разработка R&D интерфейса, внедрение Tailwind CSS стилей, интеграция сайдбара.\n"
        "• ИТОГИ:\n"
        "  - Развернут рабочий R&D стенд: danil.goida2.ru/graph\n"
        "  - Код открыт в репозитории: github.com/DanilMini1/graph-visualizer",
        13, color=RGBColor(51, 51, 51)
    )
    
    prs.save(output_path)
    print(f"Successfully populated and saved presentation to {output_path}")

if __name__ == "__main__":
    main()
